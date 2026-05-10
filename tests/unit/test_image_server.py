from __future__ import annotations

import base64
import io
import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

_MINIMAL_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk"
    "+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)


def _load_image_server(tmp_path: Path, monkeypatch):
    """Import image_server with ROOT set to tmp_path."""
    monkeypatch.setattr(sys, "argv", ["crisai-test-vision", str(tmp_path)])
    for name in list(sys.modules):
        if name == "crisai.servers.image_server" or name.startswith("crisai.servers.image_server."):
            del sys.modules[name]
    import crisai.servers.image_server as mod
    return mod


def _pptx_bytes_with_images(slide_image_counts: list[int]) -> bytes:
    prs = Presentation()
    for count in slide_image_counts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        for _ in range(count):
            slide.shapes.add_picture(
                io.BytesIO(_MINIMAL_PNG), Inches(1), Inches(1), Inches(1), Inches(1)
            )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pptx_bytes_no_images() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if slide.shapes.title:
        slide.shapes.title.text = "Text only"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# describe_image
# ---------------------------------------------------------------------------


def test_describe_image_returns_description(tmp_path, monkeypatch):
    mod = _load_image_server(tmp_path, monkeypatch)
    (tmp_path / "photo.png").write_bytes(_MINIMAL_PNG)
    monkeypatch.setattr(mod, "_describe_image_blob", lambda blob, ct, p: "a small red dot")

    result = mod.describe_image("photo.png")

    assert result == "a small red dot"


def test_describe_image_rejects_unsupported_extension(tmp_path, monkeypatch):
    mod = _load_image_server(tmp_path, monkeypatch)
    (tmp_path / "file.doc").write_bytes(b"not an image")

    with pytest.raises(ValueError, match="Unsupported image type"):
        mod.describe_image("file.doc")


def test_describe_image_raises_for_missing_file(tmp_path, monkeypatch):
    mod = _load_image_server(tmp_path, monkeypatch)

    with pytest.raises(FileNotFoundError):
        mod.describe_image("nonexistent.png")


# ---------------------------------------------------------------------------
# describe_powerpoint_slide_images
# ---------------------------------------------------------------------------


def test_describe_powerpoint_slide_images_returns_per_slide_results(tmp_path, monkeypatch):
    mod = _load_image_server(tmp_path, monkeypatch)
    (tmp_path / "deck.pptx").write_bytes(_pptx_bytes_with_images([1]))
    monkeypatch.setattr(mod, "_describe_image_blob", lambda blob, ct, p: "bar chart")

    results = mod.describe_powerpoint_slide_images("deck.pptx")

    assert len(results) == 1
    assert results[0]["slide_number"] == 1
    assert results[0]["image_index"] == 0
    assert results[0]["description"] == "bar chart"
    assert "content_type" in results[0]


def test_describe_powerpoint_slide_images_filters_slide_numbers(tmp_path, monkeypatch):
    mod = _load_image_server(tmp_path, monkeypatch)
    (tmp_path / "deck.pptx").write_bytes(_pptx_bytes_with_images([1, 1]))
    monkeypatch.setattr(mod, "_describe_image_blob", lambda blob, ct, p: "image")

    results = mod.describe_powerpoint_slide_images("deck.pptx", slide_numbers=[2])

    assert len(results) == 1
    assert results[0]["slide_number"] == 2


def test_describe_powerpoint_slide_images_empty_for_no_images(tmp_path, monkeypatch):
    mod = _load_image_server(tmp_path, monkeypatch)
    (tmp_path / "text_only.pptx").write_bytes(_pptx_bytes_no_images())
    monkeypatch.setattr(mod, "_describe_image_blob", lambda blob, ct, p: "never called")

    results = mod.describe_powerpoint_slide_images("text_only.pptx")

    assert results == []


def test_describe_powerpoint_slide_images_rejects_non_pptx(tmp_path, monkeypatch):
    mod = _load_image_server(tmp_path, monkeypatch)
    (tmp_path / "report.pdf").write_bytes(b"%PDF-1.4")

    with pytest.raises(ValueError, match="Expected a .pptx file"):
        mod.describe_powerpoint_slide_images("report.pdf")
