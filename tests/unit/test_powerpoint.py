from __future__ import annotations

import base64
import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from crisai.powerpoint import (
    extract_powerpoint_from_bytes,
    extract_powerpoint_from_path,
    extract_slide_images,
)

# Minimal valid 1x1 pixel PNG for embedding in test slides.
_MINIMAL_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk"
    "+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)


def _pptx_with_images(*slide_image_counts: int) -> bytes:
    """Build a PPTX where each positional arg specifies how many images to add per slide."""
    prs = Presentation()
    for count in slide_image_counts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        for _ in range(count):
            slide.shapes.add_picture(
                io.BytesIO(_MINIMAL_PNG), Inches(1), Inches(1), Inches(1), Inches(1)
            )
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _sample_pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Integration Strategy"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(1))
    textbox.text_frame.text = "Strategic objectives and current challenges"

    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2.5), Inches(5), Inches(1))
    table = table_shape.table
    table.cell(0, 0).text = "Theme"
    table.cell(0, 1).text = "Detail"
    table.cell(1, 0).text = "Roadmap"
    table.cell(1, 1).text = "Next steps"

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def test_extract_powerpoint_from_bytes_returns_slide_level_text_and_tables() -> None:
    extraction = extract_powerpoint_from_bytes(_sample_pptx_bytes())

    assert extraction.status == "partial_text"
    assert extraction.slide_count == 1
    assert extraction.coverage == ["slide_text", "tables"]
    assert "image_text_not_extracted" in extraction.limitations

    slide = extraction.slides[0]
    assert slide.title == "Integration Strategy"
    assert "Strategic objectives and current challenges" in slide.text
    assert slide.tables == [[["Theme", "Detail"], ["Roadmap", "Next steps"]]]


def test_extract_slide_images_returns_image_per_picture_shape() -> None:
    images = extract_slide_images(_pptx_with_images(1))

    assert len(images) == 1
    assert images[0]["slide_number"] == 1
    assert images[0]["image_index"] == 0
    assert isinstance(images[0]["blob"], bytes)
    assert len(images[0]["blob"]) > 0


def test_extract_slide_images_empty_when_no_pictures() -> None:
    assert extract_slide_images(_sample_pptx_bytes()) == []


def test_extract_slide_images_respects_slide_order() -> None:
    images = extract_slide_images(_pptx_with_images(1, 1))

    assert len(images) == 2
    assert images[0]["slide_number"] == 1
    assert images[1]["slide_number"] == 2
    assert images[0]["image_index"] == 0
    assert images[1]["image_index"] == 0


def test_extract_powerpoint_from_path_renders_extraction_header(tmp_path: Path) -> None:
    path = tmp_path / "deck.pptx"
    path.write_bytes(_sample_pptx_bytes())

    rendered = extract_powerpoint_from_path(path).to_text()

    assert "[PowerPoint extraction]" in rendered
    assert "coverage: slide_text, tables" in rendered
    assert "[Slide 1]" in rendered
    assert "Title: Integration Strategy" in rendered
    assert "Roadmap | Next steps" in rendered
