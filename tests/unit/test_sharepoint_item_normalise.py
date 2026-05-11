from __future__ import annotations

import importlib
import io
import sys

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation


def _load_sharepoint_server(monkeypatch: pytest.MonkeyPatch, tmp_path):
    monkeypatch.setattr(sys, "argv", ["sharepoint_server.py", str(tmp_path)])
    import crisai.servers.sharepoint_server as sharepoint_server

    return importlib.reload(sharepoint_server)


def _sample_pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Integration Strategy"
    textbox = slide.shapes.add_textbox(914400, 1371600, 4572000, 914400)
    textbox.text_frame.text = "Target architecture and roadmap"
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _sample_docx_bytes() -> bytes:
    doc = Document()
    doc.add_paragraph("Document paragraph")
    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()


def _sample_xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["Name", "Value"])
    sheet.append(["Alpha", 1])
    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue()


def test_normalise_item_duplicates_web_url_as_open_url(monkeypatch: pytest.MonkeyPatch, tmp_path):
    sharepoint_server = _load_sharepoint_server(monkeypatch, tmp_path)
    row = sharepoint_server._normalise_item(
        {
            "id": "item-1",
            "name": "Plan.pdf",
            "webUrl": "https://contoso.sharepoint.com/sites/a/Shared%20Documents/Plan.pdf",
            "parentReference": {"driveId": "drive-1"},
            "file": {"mimeType": "application/pdf"},
        }
    )
    assert row["webUrl"] == "https://contoso.sharepoint.com/sites/a/Shared%20Documents/Plan.pdf"
    assert row["open_url"] == row["webUrl"]
    assert row["read_handle"].startswith("sharepoint_doc:")
    assert sharepoint_server._decode_read_handle(row["read_handle"]) == ("drive-1", "item-1")


def test_decode_read_handle_rejects_malformed_handle(monkeypatch: pytest.MonkeyPatch, tmp_path):
    sharepoint_server = _load_sharepoint_server(monkeypatch, tmp_path)
    with pytest.raises(ValueError, match="read_handle"):
        sharepoint_server._decode_read_handle("not-a-handle")


def test_read_sharepoint_document_by_handle_delegates_to_raw_reader(monkeypatch: pytest.MonkeyPatch, tmp_path):
    sharepoint_server = _load_sharepoint_server(monkeypatch, tmp_path)
    handle = sharepoint_server._encode_read_handle("drive-1", "item-1")
    calls: list[tuple[str, str]] = []

    def fake_read_sharepoint_document(drive_id: str, item_id: str) -> str:
        calls.append((drive_id, item_id))
        return "deck text"

    monkeypatch.setattr(sharepoint_server, "read_sharepoint_document", fake_read_sharepoint_document)

    assert sharepoint_server.read_sharepoint_document_by_handle(handle) == "deck text"
    assert calls == [("drive-1", "item-1")]


def test_inspect_sharepoint_powerpoint_by_handle_returns_slide_extraction(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path,
):
    sharepoint_server = _load_sharepoint_server(monkeypatch, tmp_path)
    handle = sharepoint_server._encode_read_handle("drive-1", "item-1")

    monkeypatch.setattr(
        sharepoint_server,
        "_graph_get",
        lambda path: {
            "id": "item-1",
            "name": "Deck.pptx",
            "webUrl": "https://example.com/Deck.pptx",
            "parentReference": {"driveId": "drive-1"},
            "file": {"mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"},
        },
    )
    monkeypatch.setattr(sharepoint_server, "_graph_get_bytes", lambda url: _sample_pptx_bytes())

    result = sharepoint_server.inspect_sharepoint_powerpoint_by_handle(handle)

    assert result["status"] == "partial_text"
    assert result["slide_count"] == 1
    assert result["source"]["read_handle"] == handle
    assert result["slides"][0]["title"] == "Integration Strategy"
    assert "Target architecture and roadmap" in result["slides"][0]["text"]


def test_extract_bytes_by_suffix_supports_common_office_and_text_formats(monkeypatch: pytest.MonkeyPatch, tmp_path):
    sharepoint_server = _load_sharepoint_server(monkeypatch, tmp_path)

    assert sharepoint_server._extract_bytes_by_suffix(b"hello", ".txt") == "hello"
    assert "a | b" in sharepoint_server._extract_bytes_by_suffix(b"a,b\n1,2\n", ".csv")
    assert "Document paragraph" in sharepoint_server._extract_bytes_by_suffix(_sample_docx_bytes(), ".docx")
    assert "Alpha | 1" in sharepoint_server._extract_bytes_by_suffix(_sample_xlsx_bytes(), ".xlsx")
    with pytest.raises(ValueError, match="Unsupported"):
        sharepoint_server._extract_bytes_by_suffix(b"binary", ".bin")


def test_read_sharepoint_document_downloads_and_extracts_by_suffix(monkeypatch: pytest.MonkeyPatch, tmp_path):
    sharepoint_server = _load_sharepoint_server(monkeypatch, tmp_path)

    monkeypatch.setattr(
        sharepoint_server,
        "_graph_get",
        lambda path: {"id": "item-1", "name": "notes.txt", "parentReference": {"driveId": "drive-1"}},
    )
    monkeypatch.setattr(sharepoint_server, "_graph_get_bytes", lambda url: b"plain text")

    assert sharepoint_server.read_sharepoint_document("drive-1", "item-1") == "plain text"


def test_list_and_search_tools_normalise_graph_results(monkeypatch: pytest.MonkeyPatch, tmp_path):
    sharepoint_server = _load_sharepoint_server(monkeypatch, tmp_path)
    calls: list[tuple[str, dict | None, int]] = []

    def fake_graph_get(path, params=None, timeout=60):
        calls.append((path, params, timeout))
        if path == "/sites":
            return {
                "value": [
                    {
                        "id": "site-1",
                        "name": "Architecture",
                        "displayName": "Architecture",
                        "webUrl": "https://contoso.sharepoint.com/sites/architecture",
                        "siteCollection": {"hostname": "contoso.sharepoint.com"},
                    }
                ]
            }
        return {
            "value": [
                {
                    "id": "item-1",
                    "name": "Strategy.docx",
                    "webUrl": "https://example/Strategy.docx",
                    "parentReference": {"driveId": "drive-1"},
                    "file": {"mimeType": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
                }
            ]
        }

    monkeypatch.setattr(sharepoint_server, "_graph_get", fake_graph_get)

    sites = sharepoint_server.list_sites(query="Architecture", max_hits=1)
    drive_hits = sharepoint_server.search_drive_documents("drive-1", "Strategy", max_hits=1)
    site_hits = sharepoint_server.search_site_drive_documents("site-1", "Strategy", max_hits=1)

    assert sites[0]["open_url"] == "https://contoso.sharepoint.com/sites/architecture"
    assert drive_hits[0]["read_handle"].startswith("sharepoint_doc:")
    assert site_hits[0]["open_url"] == "https://example/Strategy.docx"
    assert calls[0] == ("/sites", {"search": "Architecture"}, 90)


def test_search_sharepoint_site_documents_skips_personal_sites(monkeypatch: pytest.MonkeyPatch, tmp_path):
    sharepoint_server = _load_sharepoint_server(monkeypatch, tmp_path)
    searched_sites: list[str] = []

    monkeypatch.setattr(
        sharepoint_server,
        "list_sites",
        lambda query, max_hits: [
            {"id": "personal", "displayName": "OneDrive", "webUrl": "https://contoso-my.sharepoint.com/personal/user"},
            {"id": "team", "displayName": "Team", "webUrl": "https://contoso.sharepoint.com/sites/team"},
        ],
    )

    def fake_search_site_drive_documents(site_id: str, query: str, max_hits: int):
        searched_sites.append(site_id)
        return [{"id": "item-1", "name": "Strategy.docx"}]

    monkeypatch.setattr(sharepoint_server, "search_site_drive_documents", fake_search_site_drive_documents)

    result = sharepoint_server.search_sharepoint_site_documents("Strategy")

    assert searched_sites == ["team"]
    assert result[0]["site_display_name"] == "Team"


def test_sharepoint_auth_tools_delegate_to_ms_graph(monkeypatch: pytest.MonkeyPatch, tmp_path):
    sharepoint_server = _load_sharepoint_server(monkeypatch, tmp_path)
    forced: list[bool] = []

    monkeypatch.setattr(sharepoint_server.ms_graph, "acquire_token", lambda force_interactive=False: forced.append(force_interactive))
    monkeypatch.setattr(sharepoint_server.ms_graph, "read_token_info", lambda: {"account": "me", "scope": "Files.Read"})
    monkeypatch.setattr(sharepoint_server.ms_graph, "delegated_auth_status", lambda: {"authenticated": True})
    monkeypatch.setattr(sharepoint_server, "_graph_get", lambda path: {"displayName": "Me"})

    assert "Account=me" in sharepoint_server.login_sharepoint()
    assert forced == [True]
    assert sharepoint_server.sharepoint_auth_status() == {"authenticated": True}
    assert sharepoint_server.who_am_i() == {"displayName": "Me"}
