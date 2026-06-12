from __future__ import annotations

import importlib.util
import io
import json
import sys
from pathlib import Path
from types import ModuleType

import pytest
from PIL import Image
from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parents[2]
DOCUMENT_SERVER_CANDIDATES = [
    REPO_ROOT / "src" / "crisai" / "servers" / "document_server.py",
    REPO_ROOT / "src" / "crisai" / "mcp" / "document_server.py",
    REPO_ROOT / "servers" / "document_server.py",
    REPO_ROOT / "document_server.py",
]


def _find_document_server() -> Path:
    """Find the document MCP server without assuming one exact repo layout."""
    for candidate in DOCUMENT_SERVER_CANDIDATES:
        if candidate.exists():
            return candidate

    searched = "\n".join(str(path) for path in DOCUMENT_SERVER_CANDIDATES)
    raise FileNotFoundError(f"Could not find document_server.py. Searched:\n{searched}")


@pytest.fixture()
def document_server(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> ModuleType:
    """Load document_server.py with ROOT pointing at a temporary workspace."""
    server_path = _find_document_server()
    module_name = f"document_server_under_test_{id(tmp_path)}"

    monkeypatch.setattr(sys, "argv", [str(server_path), str(tmp_path)])

    spec = importlib.util.spec_from_file_location(module_name, server_path)
    assert spec is not None
    assert spec.loader is not None

    module = importlib.util.module_from_spec(spec)
    sys.modules[module_name] = module
    spec.loader.exec_module(module)
    return module


def _write_context_file(workspace: Path, relative_path: str, content: str) -> None:
    path = workspace / "knowledge" / relative_path
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def _write_sample_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Integration Strategy"
    textbox = slide.shapes.add_textbox(914400, 1371600, 4572000, 914400)
    textbox.text_frame.text = "Strategic themes and target operating model"
    output = io.BytesIO()
    prs.save(output)
    path.write_bytes(output.getvalue())


def test_build_context_index_creates_chunks_across_context_folders(
    document_server: ModuleType,
    tmp_path: Path,
) -> None:
    _write_context_file(
        tmp_path,
        "standards/source-control.txt",
        "Recurring reports must define ownership, lineage, access, and data quality checks. "
        "Critical transformation logic must not live only in Power BI.",
    )
    _write_context_file(
        tmp_path,
        "patterns/file-route.txt",
        "For recurring Excel reporting, land the file unchanged, validate fields, stage the data, "
        "curate it, and expose Power BI from the governed dataset.",
    )
    _write_context_file(
        tmp_path,
        "notes/session-capture.txt",
        "Stakeholders want a quick dashboard from spreadsheets, but they expect monthly refreshes "
        "and consistent figures for operational decisions.",
    )

    summary = document_server.build_context_index(
        context_subdir="knowledge",
        max_chars=220,
        overlap_chars=40,
    )

    assert summary["documents_indexed"] == 3
    assert summary["chunks_indexed"] >= 3
    assert set(summary["folder_counts"]).issuperset({"standards", "patterns", "notes"})
    assert (tmp_path / ".crisai" / "context_index.json").exists()

    index_summary = document_server.get_context_index_summary()

    assert index_summary["context_subdir"] == "knowledge"
    assert index_summary["documents_indexed"] == 3
    assert index_summary["chunking"] == {"max_chars": 220, "overlap_chars": 40}


def test_search_context_chunks_returns_ranked_cross_folder_results(
    document_server: ModuleType,
    tmp_path: Path,
) -> None:
    _write_context_file(
        tmp_path,
        "reference/background.txt",
        "Local Excel files often become operational reporting sources. The architecture must clarify "
        "whether the dashboard is ad hoc analysis or a recurring reporting product.",
    )
    _write_context_file(
        tmp_path,
        "standards/assurance.txt",
        "Recurring reports require documented ownership, lineage, access model, retention expectations, "
        "and known data quality limitations.",
    )
    _write_context_file(
        tmp_path,
        "patterns/controlled-ingestion.txt",
        "The preferred pattern for recurring Excel dashboards is controlled ingestion: land the file, "
        "validate structure, curate a dataset, then present the result in Power BI.",
    )
    _write_context_file(
        tmp_path,
        "designs/previous-example.txt",
        "A previous benefits dashboard used Excel as a source but ingested the file into the data platform "
        "before Power BI because monthly refresh and assurance were required.",
    )
    _write_context_file(
        tmp_path,
        "notes/current-discovery.txt",
        "The current team wants speed, but the report will refresh monthly and influence operational decisions. "
        "The source spreadsheet has no formal owner yet.",
    )

    results = document_server.search_context_chunks(
        "recurring Excel Power BI dashboard ownership lineage data quality monthly refresh",
        max_results=8,
        rebuild=True,
        context_subdir="knowledge",
    )

    assert results
    assert results == sorted(results, key=lambda item: item["score"], reverse=True)

    folders = {result["folder"] for result in results}
    assert len(folders) >= 3
    assert {"standards", "patterns"}.issubset(folders)

    combined_text = "\n".join(result["text"].lower() for result in results)
    assert "ownership" in combined_text
    assert "lineage" in combined_text
    assert "power bi" in combined_text

    for result in results:
        assert result["score"] > 0
        assert result["path"].startswith("knowledge/")
        assert result["chunk_id"].startswith("chunk-")


def test_context_search_returns_empty_list_for_non_matching_query(
    document_server: ModuleType,
    tmp_path: Path,
) -> None:
    _write_context_file(
        tmp_path,
        "patterns/reporting.txt",
        "Power BI reporting over curated datasets should avoid hidden transformation logic.",
    )

    results = document_server.search_context_chunks(
        "kubernetes gpu model serving ollama condenser",
        max_results=5,
        rebuild=True,
        context_subdir="knowledge",
    )

    assert results == []


def test_chunk_text_validates_chunk_settings(document_server: ModuleType) -> None:
    with pytest.raises(ValueError, match="max_chars"):
        document_server._chunk_text("example", max_chars=0, overlap_chars=0)

    with pytest.raises(ValueError, match="overlap_chars"):
        document_server._chunk_text("example", max_chars=100, overlap_chars=100)


def test_inspect_powerpoint_document_returns_structured_slides(
    document_server: ModuleType,
    tmp_path: Path,
) -> None:
    path = tmp_path / "knowledge" / "deck.pptx"
    path.parent.mkdir(parents=True, exist_ok=True)
    _write_sample_pptx(path)

    result = document_server.inspect_powerpoint_document("knowledge/deck.pptx")

    assert result["status"] == "partial_text"
    assert result["path"] == "knowledge/deck.pptx"
    assert result["slides"][0]["title"] == "Integration Strategy"
    assert "Strategic themes and target operating model" in result["slides"][0]["text"]


def test_read_document_blocks_sensitive_secret_folder(document_server: ModuleType, tmp_path: Path) -> None:
    secret = tmp_path / ".secrets" / "secret.md"
    secret.parent.mkdir()
    secret.write_text("private", encoding="utf-8")

    with pytest.raises(ValueError, match="restricted"):
        document_server.read_document(".secrets/secret.md")


def test_list_supported_documents_omits_sensitive_dirs(
    document_server: ModuleType,
    tmp_path: Path,
) -> None:
    public = tmp_path / "knowledge" / "standard.md"
    public.parent.mkdir(parents=True)
    public.write_text("# Standard\n", encoding="utf-8")
    auth = tmp_path / ".auth" / "token.md"
    auth.parent.mkdir()
    auth.write_text("token", encoding="utf-8")
    cache = tmp_path / ".cache" / "cached.md"
    cache.parent.mkdir()
    cache.write_text("cache", encoding="utf-8")

    listed = document_server.list_supported_document_files(".")

    assert listed == ["knowledge/standard.md"]


def test_build_context_index_omits_sensitive_dirs(
    document_server: ModuleType,
    tmp_path: Path,
) -> None:
    public = tmp_path / "knowledge" / "standard.md"
    public.parent.mkdir(parents=True)
    public.write_text("public architecture guidance", encoding="utf-8")
    auth = tmp_path / ".auth" / "token.md"
    auth.parent.mkdir()
    auth.write_text("secret architecture guidance", encoding="utf-8")

    summary = document_server.build_context_index(context_subdir=".")

    assert summary["documents_indexed"] == 1
    assert summary["folder_counts"] == {"knowledge": 1}


def test_search_context_chunks_filters_stale_sensitive_index_entries(
    document_server: ModuleType,
    tmp_path: Path,
) -> None:
    index_path = tmp_path / ".crisai" / "context_index.json"
    index_path.parent.mkdir()
    normal_text = "needle public architecture guidance"
    secret_text = "needle secret token guidance"
    index_path.write_text(
        json.dumps(
            {
                "version": 1,
                "created_at": "2026-05-18T00:00:00",
                "context_subdir": ".",
                "documents": ["knowledge/public.md", ".auth/token.md"],
                "folder_counts": {"knowledge": 1, ".auth": 1},
                "chunks": [
                    {
                        "chunk_id": "chunk-public",
                        "path": "knowledge/public.md",
                        "folder": "knowledge",
                        "authority_weight": 1.0,
                        "chunk_index": 1,
                        "text": normal_text,
                        "vector": document_server._vectorise(normal_text),
                    },
                    {
                        "chunk_id": "chunk-secret",
                        "path": ".auth/token.md",
                        "folder": ".auth",
                        "authority_weight": 1.0,
                        "chunk_index": 1,
                        "text": secret_text,
                        "vector": document_server._vectorise(secret_text),
                    },
                ],
            }
        ),
        encoding="utf-8",
    )

    results = document_server.search_context_chunks("needle guidance", rebuild=False, context_subdir=".")

    assert [result["path"] for result in results] == ["knowledge/public.md"]


def _write_image_only_pdf(path: Path, pages: int = 1) -> None:
    """Write a PDF whose pages are raster images with no text layer."""
    images = [Image.new("RGB", (320, 200), "white") for _ in range(pages)]
    images[0].save(path, "PDF", save_all=True, append_images=images[1:])


def test_image_only_pdf_has_no_text_layer(tmp_path: Path) -> None:
    """Sanity: the generated fixture really has no extractable text."""
    from pypdf import PdfReader

    pdf_path = tmp_path / "scan.pdf"
    _write_image_only_pdf(pdf_path)
    reader = PdfReader(str(pdf_path))
    assert all(not (page.extract_text() or "").strip() for page in reader.pages)


def test_read_pdf_uses_vision_for_image_only_pages(
    document_server: ModuleType, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.delenv("CRISAI_PDF_VISION_MAX_PAGES", raising=False)
    monkeypatch.setattr(document_server, "describe_image_blob", lambda blob, ct, prompt: "RECOVERED PAGE TEXT")

    pdf_path = tmp_path / "scan.pdf"
    _write_image_only_pdf(pdf_path)

    result = document_server._read_pdf(pdf_path)

    assert "[Page 1 (vision)]" in result
    assert "RECOVERED PAGE TEXT" in result


def test_read_pdf_vision_disabled_returns_note(
    document_server: ModuleType, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setenv("CRISAI_PDF_VISION_MAX_PAGES", "0")

    def _must_not_call(*_a: object, **_kw: object) -> str:
        raise AssertionError("vision model must not be called when the cap is 0")

    monkeypatch.setattr(document_server, "describe_image_blob", _must_not_call)

    pdf_path = tmp_path / "scan.pdf"
    _write_image_only_pdf(pdf_path)

    result = document_server._read_pdf(pdf_path)

    assert "image-only page(s) with no text layer" in result
    assert "CRISAI_PDF_VISION_MAX_PAGES" in result


def test_read_pdf_respects_page_cap(
    document_server: ModuleType, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setenv("CRISAI_PDF_VISION_MAX_PAGES", "1")
    calls: list[int] = []

    def _describe(blob: bytes, ct: str, prompt: str) -> str:
        calls.append(1)
        return "PAGE TEXT"

    monkeypatch.setattr(document_server, "describe_image_blob", _describe)

    pdf_path = tmp_path / "scan.pdf"
    _write_image_only_pdf(pdf_path, pages=3)

    result = document_server._read_pdf(pdf_path)

    assert len(calls) == 1  # only one page described, cap respected
    assert "[Page 1 (vision)]" in result
    assert "vision page cap" in result
