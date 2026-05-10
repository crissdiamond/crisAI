"""PowerPoint text extraction helpers shared by MCP document tools."""

from __future__ import annotations

import io
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from pptx import Presentation

TEXT_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
REL_NS = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}
REL_TYPE_NOTES = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
IGNORED_NOTE_TEXT = {
    "click to add notes",
    "notes",
}


@dataclass(frozen=True, slots=True)
class PowerPointSlide:
    """Text extracted from one PowerPoint slide."""

    number: int
    title: str = ""
    text: list[str] = field(default_factory=list)
    tables: list[list[list[str]]] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "slide_number": self.number,
            "title": self.title,
            "text": list(self.text),
            "tables": [[list(row) for row in table] for table in self.tables],
            "speaker_notes": list(self.notes),
        }


@dataclass(frozen=True, slots=True)
class PowerPointExtraction:
    """Structured PowerPoint extraction result with coverage metadata."""

    slide_count: int
    slides: list[PowerPointSlide]
    coverage: list[str]
    limitations: list[str]

    @property
    def status(self) -> str:
        if not self.slides:
            return "no_text"
        return "partial_text"

    def to_dict(self) -> dict[str, Any]:
        return {
            "status": self.status,
            "slide_count": self.slide_count,
            "slides_with_text": len(self.slides),
            "coverage": list(self.coverage),
            "limitations": list(self.limitations),
            "slides": [slide.to_dict() for slide in self.slides],
        }

    def to_text(self) -> str:
        header = (
            "[PowerPoint extraction]\n"
            f"status: {self.status}\n"
            f"slides: {self.slide_count}\n"
            f"slides_with_text: {len(self.slides)}\n"
            f"coverage: {', '.join(self.coverage) or 'none'}\n"
            f"limitations: {', '.join(self.limitations) or 'none'}"
        )
        rendered_slides = [slide_to_text(slide) for slide in self.slides]
        return "\n\n".join([header, *rendered_slides]).strip()


def extract_powerpoint_from_path(file_path: Path) -> PowerPointExtraction:
    """Extract structured text from a PowerPoint file."""
    return extract_powerpoint_from_bytes(file_path.read_bytes())


def extract_powerpoint_from_bytes(data: bytes) -> PowerPointExtraction:
    """Extract structured text, table cells, and notes from PowerPoint bytes."""
    prs = Presentation(io.BytesIO(data))
    notes_by_slide = _extract_notes_by_slide_number(data)
    slides: list[PowerPointSlide] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = _normalise_text(slide.shapes.title.text) if slide.shapes.title is not None else ""
        text_parts: list[str] = []
        tables: list[list[list[str]]] = []
        _extract_shapes(slide.shapes, text_parts=text_parts, tables=tables)
        text_parts = _dedupe_preserve_order(part for part in text_parts if part != title)
        notes = notes_by_slide.get(idx, [])
        if title or text_parts or tables or notes:
            slides.append(PowerPointSlide(number=idx, title=title, text=text_parts, tables=tables, notes=notes))

    coverage = ["slide_text", "tables"]
    if any(slide.notes for slide in slides):
        coverage.append("speaker_notes")
    limitations = [
        "image_text_not_extracted",
        "embedded_objects_not_extracted",
        "smartart_text_may_be_partial",
    ]
    if "speaker_notes" not in coverage:
        limitations.append("speaker_notes_not_found")

    return PowerPointExtraction(
        slide_count=len(prs.slides),
        slides=slides,
        coverage=coverage,
        limitations=limitations,
    )


def slide_to_text(slide: PowerPointSlide) -> str:
    """Render one structured slide extraction as readable text."""
    lines = [f"[Slide {slide.number}]"]
    if slide.title:
        lines.append(f"Title: {slide.title}")
    lines.extend(slide.text)
    for table_index, table in enumerate(slide.tables, start=1):
        lines.append(f"Table {table_index}:")
        lines.extend(" | ".join(row) for row in table if any(cell.strip() for cell in row))
    if slide.notes:
        lines.append("Speaker notes:")
        lines.extend(slide.notes)
    return "\n".join(lines)


def _extract_shapes(shapes: Any, *, text_parts: list[str], tables: list[list[list[str]]]) -> None:
    for shape in shapes:
        if getattr(shape, "has_table", False):
            table = [
                [_normalise_text(cell.text) for cell in row.cells]
                for row in shape.table.rows
            ]
            if any(any(cell for cell in row) for row in table):
                tables.append(table)
        text = _normalise_text(getattr(shape, "text", ""))
        if text:
            text_parts.append(text)
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            _extract_shapes(child_shapes, text_parts=text_parts, tables=tables)


def _extract_notes_by_slide_number(data: bytes) -> dict[int, list[str]]:
    notes: dict[int, list[str]] = {}
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        for slide_number in _slide_numbers(archive):
            rel_path = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
            if rel_path not in archive.namelist():
                continue
            notes_path = _notes_path_from_rels(archive.read(rel_path))
            if notes_path and notes_path in archive.namelist():
                extracted = _extract_text_from_xml(archive.read(notes_path))
                if extracted:
                    notes[slide_number] = extracted
    return notes


def _slide_numbers(archive: zipfile.ZipFile) -> list[int]:
    numbers: list[int] = []
    prefix = "ppt/slides/slide"
    for name in archive.namelist():
        if not name.startswith(prefix) or not name.endswith(".xml"):
            continue
        raw = name[len(prefix) : -len(".xml")]
        if raw.isdigit():
            numbers.append(int(raw))
    return sorted(numbers)


def _notes_path_from_rels(data: bytes) -> str:
    root = ElementTree.fromstring(data)
    for relationship in root.findall("rel:Relationship", REL_NS):
        if relationship.attrib.get("Type") != REL_TYPE_NOTES:
            continue
        target = relationship.attrib.get("Target", "")
        if target.startswith("../"):
            target = target[3:]
        if target.startswith("notesSlides/"):
            return "ppt/" + target
    return ""


def _extract_text_from_xml(data: bytes) -> list[str]:
    root = ElementTree.fromstring(data)
    values = [_normalise_text(node.text or "") for node in root.findall(".//a:t", TEXT_NS)]
    return [
        value
        for value in _dedupe_preserve_order(values)
        if value and value.lower() not in IGNORED_NOTE_TEXT
    ]


def _normalise_text(text: str) -> str:
    return " ".join((text or "").split())


def _dedupe_preserve_order(values: Any) -> list[str]:
    seen: set[str] = set()
    out: list[str] = []
    for value in values:
        if value and value not in seen:
            out.append(value)
            seen.add(value)
    return out
