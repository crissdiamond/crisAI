from __future__ import annotations

import csv
import hashlib
import io
import json
import logging
import math
import re
import sys
from collections import Counter
from datetime import datetime
from pathlib import Path
from typing import Any

import chardet
from docx import Document as DocxDocument
from mcp.server.fastmcp import FastMCP
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from crisai.config import load_settings


mcp = FastMCP("crisai-document-reader")

ROOT = Path(sys.argv[1]).resolve() if len(sys.argv) > 1 else Path.cwd().resolve()
ROOT.mkdir(parents=True, exist_ok=True)

LOG_FILE = load_settings().log_dir / "document_mcp.log"
CONTEXT_INDEX_FILE = ROOT / ".crisai" / "context_index.json"
SUPPORTED_DOCUMENT_SUFFIXES = {
    ".txt",
    ".md",
    ".json",
    ".yaml",
    ".yml",
    ".py",
    ".log",
    ".csv",
    ".docx",
    ".pdf",
    ".pptx",
    ".xlsx",
}
CONTEXT_AUTHORITY_WEIGHTS = {
    "standards": 1.30,
    "patterns": 1.20,
    "designs": 1.10,
    "reference": 1.00,
    "notes": 0.90,
}
TOKEN_PATTERN = re.compile(r"[a-zA-Z][a-zA-Z0-9_-]{1,}")
STOP_WORDS = {
    "about",
    "after",
    "again",
    "against",
    "also",
    "because",
    "before",
    "being",
    "between",
    "could",
    "from",
    "have",
    "into",
    "more",
    "must",
    "only",
    "should",
    "than",
    "that",
    "their",
    "then",
    "there",
    "these",
    "this",
    "through",
    "used",
    "when",
    "where",
    "which",
    "with",
    "would",
}


def _configure_mcp_logging() -> None:
    """Keep MCP framework INFO logs out of the interactive CLI.

    Warnings and errors are still written to this server log file.
    """
    logger_names = [
        "mcp",
        "mcp.server",
        "mcp.server.fastmcp",
        "mcp.server.lowlevel",
        "mcp.server.lowlevel.server",
    ]

    formatter = logging.Formatter("%(asctime)s %(levelname)s %(name)s %(message)s")
    file_handler = logging.FileHandler(LOG_FILE, encoding="utf-8")
    file_handler.setLevel(logging.WARNING)
    file_handler.setFormatter(formatter)

    for name in logger_names:
        logger = logging.getLogger(name)
        logger.setLevel(logging.WARNING)
        logger.handlers.clear()
        logger.addHandler(file_handler)
        logger.propagate = False


def log_event(message: str) -> None:
    with LOG_FILE.open("a", encoding="utf-8") as f:
        f.write(f"{datetime.now().isoformat()} {message}\n")


_configure_mcp_logging()

log_event(f"server_started root={ROOT}")


def _safe_path(relative_path: str) -> Path:
    raw = (relative_path or ".").strip()

    if raw.startswith("/"):
        raw = raw.lstrip("/")
    if raw.startswith("workspace/"):
        raw = raw[len("workspace/"):]

    candidate = (ROOT / raw).resolve()
    root = ROOT.resolve()

    if candidate != root and root not in candidate.parents:
        raise ValueError(
            f"Path escapes the workspace root. root={root} requested={relative_path} resolved={candidate}"
        )

    return candidate


def _detect_text_encoding(data: bytes) -> str:
    detected = chardet.detect(data)
    encoding = detected.get("encoding") or "utf-8"
    return encoding


def _read_text_like(file_path: Path) -> str:
    data = file_path.read_bytes()
    encoding = _detect_text_encoding(data)
    try:
        return data.decode(encoding)
    except Exception:
        return data.decode("utf-8", errors="replace")


def _read_docx(file_path: Path) -> str:
    doc = DocxDocument(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n\n".join(paragraphs)


def _read_pdf(file_path: Path) -> str:
    reader = PdfReader(str(file_path))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            pages.append(f"[Page {i}]\n{text}")
    return "\n\n".join(pages)


def _read_pptx(file_path: Path) -> str:
    prs = Presentation(str(file_path))
    slides_out = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    parts.append(text)
        if parts:
            slides_out.append(f"[Slide {idx}]\n" + "\n".join(parts))
    return "\n\n".join(slides_out)


def _read_xlsx(file_path: Path, max_rows: int = 50, max_cols: int = 20) -> str:
    wb = load_workbook(filename=str(file_path), read_only=True, data_only=True)
    out = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        out.append(f"[Sheet: {sheet_name}]")
        row_count = 0

        for row in ws.iter_rows(values_only=True):
            if row_count >= max_rows:
                out.append("... truncated ...")
                break

            values = []
            for cell in row[:max_cols]:
                if cell is None:
                    values.append("")
                else:
                    values.append(str(cell))

            if any(v != "" for v in values):
                out.append(" | ".join(values))
                row_count += 1

        out.append("")

    return "\n".join(out).strip()


def _read_csv(file_path: Path, max_rows: int = 100) -> str:
    data = file_path.read_bytes()
    encoding = _detect_text_encoding(data)
    text = data.decode(encoding, errors="replace")

    reader = csv.reader(io.StringIO(text))
    rows = []
    for idx, row in enumerate(reader):
        if idx >= max_rows:
            rows.append(["... truncated ..."])
            break
        rows.append(row)

    return "\n".join(" | ".join(cell for cell in row) for row in rows)


def _read_supported_document(file_path: Path) -> str:
    """Extract text from any supported document type."""
    suffix = file_path.suffix.lower()

    if suffix in {".txt", ".md", ".json", ".yaml", ".yml", ".py", ".log"}:
        return _read_text_like(file_path)

    if suffix == ".csv":
        return _read_csv(file_path)

    if suffix == ".docx":
        return _read_docx(file_path)

    if suffix == ".pdf":
        return _read_pdf(file_path)

    if suffix == ".pptx":
        return _read_pptx(file_path)

    if suffix == ".xlsx":
        return _read_xlsx(file_path)

    raise ValueError(f"Unsupported document type: {suffix or 'no extension'}")


def _tokenize(text: str) -> list[str]:
    """Return normalised tokens for lightweight local vectorisation."""
    tokens = [token.lower() for token in TOKEN_PATTERN.findall(text)]
    return [token for token in tokens if token not in STOP_WORDS]


def _vectorise(text: str) -> dict[str, float]:
    """Create a sparse term-frequency vector without external dependencies."""
    tokens = _tokenize(text)
    if not tokens:
        return {}

    counts = Counter(tokens)
    norm = math.sqrt(sum(value * value for value in counts.values()))
    if norm == 0:
        return {}

    return {token: value / norm for token, value in counts.items()}


def _cosine_similarity(left: dict[str, float], right: dict[str, float]) -> float:
    """Calculate cosine similarity for two sparse vectors."""
    if not left or not right:
        return 0.0

    smaller, larger = (left, right) if len(left) <= len(right) else (right, left)
    return sum(value * larger.get(token, 0.0) for token, value in smaller.items())


def _stable_chunk_id(path: str, chunk_index: int, text: str) -> str:
    """Build a stable chunk identifier from path, index, and content."""
    digest = hashlib.sha1(f"{path}:{chunk_index}:{text}".encode("utf-8")).hexdigest()[:12]
    return f"chunk-{digest}"


def _chunk_text(text: str, max_chars: int = 1200, overlap_chars: int = 200) -> list[str]:
    """Split text into overlapping chunks while preserving paragraph boundaries.

    Args:
        text: Extracted document text.
        max_chars: Preferred maximum size for a chunk.
        overlap_chars: Approximate character overlap between adjacent chunks.

    Returns:
        A list of non-empty text chunks.
    """
    if max_chars <= 0:
        raise ValueError("max_chars must be greater than zero")
    if overlap_chars < 0:
        raise ValueError("overlap_chars cannot be negative")
    if overlap_chars >= max_chars:
        raise ValueError("overlap_chars must be smaller than max_chars")

    paragraphs = [part.strip() for part in re.split(r"\n\s*\n", text) if part.strip()]
    chunks: list[str] = []
    current = ""

    for paragraph in paragraphs:
        if len(paragraph) > max_chars:
            if current:
                chunks.append(current.strip())
                current = ""
            start = 0
            step = max_chars - overlap_chars
            while start < len(paragraph):
                chunks.append(paragraph[start:start + max_chars].strip())
                start += step
            continue

        candidate = f"{current}\n\n{paragraph}".strip() if current else paragraph
        if len(candidate) <= max_chars:
            current = candidate
        else:
            chunks.append(current.strip())
            overlap = current[-overlap_chars:].strip() if overlap_chars else ""
            current = f"{overlap}\n\n{paragraph}".strip() if overlap else paragraph

    if current:
        chunks.append(current.strip())

    return chunks


def _context_folder_for(file_path: Path, context_root: Path) -> str:
    """Return the first folder under the context root for authority weighting."""
    relative_parts = file_path.relative_to(context_root).parts
    if not relative_parts:
        return "unknown"
    return relative_parts[0]


def _iter_supported_files(base: Path) -> list[Path]:
    """Return supported files under a base path, excluding hidden internal folders."""
    if not base.exists():
        return []

    return sorted(
        file_path
        for file_path in base.rglob("*")
        if file_path.is_file()
        and file_path.suffix.lower() in SUPPORTED_DOCUMENT_SUFFIXES
        and ".crisai" not in file_path.parts
    )


def _build_context_chunks(
    context_subdir: str = "context",
    max_chars: int = 1200,
    overlap_chars: int = 200,
) -> list[dict[str, Any]]:
    """Build context chunks from supported files in the context knowledge area."""
    context_root = _safe_path(context_subdir)
    chunks: list[dict[str, Any]] = []

    for file_path in _iter_supported_files(context_root):
        try:
            text = _read_supported_document(file_path)
        except Exception as exc:
            log_event(f"context_index_skip path={file_path} error={exc}")
            continue

        relative_path = str(file_path.relative_to(ROOT))
        folder = _context_folder_for(file_path, context_root)
        authority_weight = CONTEXT_AUTHORITY_WEIGHTS.get(folder, 1.0)

        for chunk_index, chunk_text in enumerate(_chunk_text(text, max_chars, overlap_chars), start=1):
            chunks.append(
                {
                    "chunk_id": _stable_chunk_id(relative_path, chunk_index, chunk_text),
                    "path": relative_path,
                    "folder": folder,
                    "authority_weight": authority_weight,
                    "chunk_index": chunk_index,
                    "text": chunk_text,
                    "vector": _vectorise(chunk_text),
                }
            )

    return chunks


def _write_context_index(index: dict[str, Any]) -> None:
    CONTEXT_INDEX_FILE.parent.mkdir(parents=True, exist_ok=True)
    CONTEXT_INDEX_FILE.write_text(json.dumps(index, indent=2), encoding="utf-8")


def _read_context_index() -> dict[str, Any]:
    if not CONTEXT_INDEX_FILE.exists():
        raise FileNotFoundError(
            f"Context index not found: {CONTEXT_INDEX_FILE.relative_to(ROOT)}. Run build_context_index first."
        )
    return json.loads(CONTEXT_INDEX_FILE.read_text(encoding="utf-8"))


@mcp.tool()
def read_document(path: str) -> str:
    """Read a document inside the workspace and extract text from common file formats."""
    log_event(f"read_document path={path}")
    file_path = _safe_path(path)

    if not file_path.exists():
        raise FileNotFoundError(f"No such file: {file_path}")

    return _read_supported_document(file_path)


@mcp.tool()
def get_document_metadata(path: str) -> dict[str, str | int]:
    """Return basic metadata about a document inside the workspace."""
    log_event(f"get_document_metadata path={path}")
    file_path = _safe_path(path)

    if not file_path.exists():
        raise FileNotFoundError(f"No such file: {file_path}")

    stat = file_path.stat()
    resolved = file_path.resolve()
    return {
        "path": str(file_path.relative_to(ROOT)),
        "name": file_path.name,
        "suffix": file_path.suffix.lower(),
        "size_bytes": stat.st_size,
        "file_uri": resolved.as_uri(),
    }


@mcp.tool()
def list_supported_document_files(subdir: str = ".") -> list[str]:
    """List supported readable document files under a workspace subdirectory."""
    log_event(f"list_supported_document_files subdir={subdir}")
    base = _safe_path(subdir)

    return [str(file_path.relative_to(ROOT)) for file_path in _iter_supported_files(base)]


@mcp.tool()
def build_context_index(
    context_subdir: str = "context",
    max_chars: int = 1200,
    overlap_chars: int = 200,
) -> dict[str, Any]:
    """Build a local vector index from documents under the context knowledge area.

    Args:
        context_subdir: Workspace-relative folder containing context knowledge.
        max_chars: Preferred maximum characters per chunk.
        overlap_chars: Approximate character overlap between chunks.

    Returns:
        Summary information about the generated index.
    """
    log_event(
        "build_context_index "
        f"context_subdir={context_subdir} max_chars={max_chars} overlap_chars={overlap_chars}"
    )
    context_root = _safe_path(context_subdir)
    chunks = _build_context_chunks(context_subdir, max_chars, overlap_chars)
    document_paths = sorted({chunk["path"] for chunk in chunks})
    folder_counts = Counter(chunk["folder"] for chunk in chunks)

    index = {
        "version": 1,
        "created_at": datetime.now().isoformat(),
        "context_subdir": str(context_root.relative_to(ROOT)),
        "chunking": {
            "max_chars": max_chars,
            "overlap_chars": overlap_chars,
        },
        "documents": document_paths,
        "folder_counts": dict(sorted(folder_counts.items())),
        "chunks": chunks,
    }
    _write_context_index(index)

    return {
        "index_path": str(CONTEXT_INDEX_FILE.relative_to(ROOT)),
        "documents_indexed": len(document_paths),
        "chunks_indexed": len(chunks),
        "folder_counts": dict(sorted(folder_counts.items())),
    }


@mcp.tool()
def search_context_chunks(
    query: str,
    max_results: int = 8,
    rebuild: bool = False,
    context_subdir: str = "context",
) -> list[dict[str, Any]]:
    """Search indexed context chunks using local vector similarity.

    Args:
        query: Natural language search query.
        max_results: Maximum number of chunks to return.
        rebuild: Whether to rebuild the index before searching.
        context_subdir: Workspace-relative folder containing context knowledge.

    Returns:
        Ranked context chunks with source path, folder, score, and text.
    """
    log_event(
        f"search_context_chunks query={query!r} max_results={max_results} "
        f"rebuild={rebuild} context_subdir={context_subdir}"
    )
    if max_results <= 0:
        return []

    if rebuild or not CONTEXT_INDEX_FILE.exists():
        build_context_index(context_subdir=context_subdir)

    index = _read_context_index()
    query_vector = _vectorise(query)
    ranked: list[dict[str, Any]] = []

    for chunk in index.get("chunks", []):
        vector_score = _cosine_similarity(query_vector, chunk.get("vector", {}))
        if vector_score <= 0:
            continue

        authority_weight = float(chunk.get("authority_weight", 1.0))
        final_score = vector_score * authority_weight
        ranked.append(
            {
                "score": round(final_score, 6),
                "vector_score": round(vector_score, 6),
                "authority_weight": authority_weight,
                "chunk_id": chunk["chunk_id"],
                "path": chunk["path"],
                "folder": chunk["folder"],
                "chunk_index": chunk["chunk_index"],
                "text": chunk["text"],
            }
        )

    ranked.sort(key=lambda item: item["score"], reverse=True)
    return ranked[:max_results]


@mcp.tool()
def get_context_index_summary() -> dict[str, Any]:
    """Return summary information for the current context index."""
    log_event("get_context_index_summary")
    index = _read_context_index()
    return {
        "version": index.get("version"),
        "created_at": index.get("created_at"),
        "context_subdir": index.get("context_subdir"),
        "documents_indexed": len(index.get("documents", [])),
        "chunks_indexed": len(index.get("chunks", [])),
        "folder_counts": index.get("folder_counts", {}),
        "chunking": index.get("chunking", {}),
    }


if __name__ == "__main__":
    mcp.run()
