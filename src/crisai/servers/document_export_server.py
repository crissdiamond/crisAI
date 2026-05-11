from __future__ import annotations

import re
import sys
from pathlib import Path
from typing import Any

import yaml
from docx import Document
from mcp.server.fastmcp import FastMCP
from pptx import Presentation

from crisai.config import load_settings
from crisai.logging_utils import append_json_log_line, configure_mcp_framework_logging
from crisai.workspace.spaces import load_workspace_spaces

mcp = FastMCP("crisai-document-export")

ROOT = Path(sys.argv[1]).resolve() if len(sys.argv) > 1 else Path.cwd().resolve()
ROOT.mkdir(parents=True, exist_ok=True)

ALLOWED_OUTPUT_SUFFIXES = {".docx", ".pptx"}
HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$")


def _log_file() -> Path:
    return load_settings().log_dir / "document_export_mcp.log"


def _configure_mcp_logging() -> None:
    """Keep MCP framework INFO logs out of the interactive CLI."""
    configure_mcp_framework_logging(_log_file(), service_component="document_export_mcp")


def log_event(message: str) -> None:
    append_json_log_line(
        _log_file(),
        message,
        logger_name="crisai.mcp.document_export",
        service_component="document_export_mcp",
    )


def _safe_path(relative_path: str) -> Path:
    raw = (relative_path or ".").strip()
    spaces = load_workspace_spaces()
    if raw.startswith("/"):
        raw = raw.lstrip("/")
    if raw.startswith("workspace/"):
        raw = raw[len("workspace/") :]
    if (raw == "context" or raw.startswith("context/")) and not (ROOT / raw).exists():
        raw = raw.replace("context", spaces.knowledge_root, 1)
    if (raw == "context_staging" or raw.startswith("context_staging/")) and not (ROOT / raw).exists():
        raw = raw.replace("context_staging", spaces.knowledge_staging_root, 1)

    candidate = (ROOT / raw).resolve()
    root = ROOT.resolve()
    if candidate != root and root not in candidate.parents:
        raise ValueError(
            f"Path escapes the workspace root. root={root} requested={relative_path} resolved={candidate}"
        )
    return candidate


def _relative(path: Path) -> str:
    return path.resolve().relative_to(ROOT.resolve()).as_posix()


def _enforce_output_path(path: Path, suffix: str) -> None:
    if path.suffix.lower() != suffix:
        raise ValueError(f"Output path must end with {suffix}.")
    if path.suffix.lower() not in ALLOWED_OUTPUT_SUFFIXES:
        raise ValueError(f"Unsupported output suffix: {path.suffix}.")
    rel = _relative(path)
    is_task_export = rel.startswith("tasks/") and "/exports/" in rel
    is_generic_output = rel == "outputs" or rel.startswith("outputs/")
    if not (is_task_export or is_generic_output):
        raise ValueError("Document exports must be written under tasks/<task>/exports/ or outputs/.")


def _parse_front_matter(raw: str) -> tuple[dict[str, Any], str]:
    if not raw.startswith("---"):
        return {}, raw
    close = raw.find("\n---\n", 3)
    if close == -1:
        return {}, raw
    try:
        metadata = yaml.safe_load(raw[3:close]) or {}
    except yaml.YAMLError:
        return {}, raw
    return (metadata if isinstance(metadata, dict) else {}), raw[close + 5 :]


def _parse_markdown(raw: str) -> dict[str, Any]:
    metadata, body = _parse_front_matter(raw)
    title = str(metadata.get("title") or "").strip()
    sections: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    intro: list[str] = []
    for line in body.splitlines():
        match = HEADING_RE.match(line)
        if match:
            current = {"level": len(match.group(1)), "title": match.group(2).strip(), "lines": []}
            sections.append(current)
            if not title and current["level"] == 1:
                title = str(current["title"])
            continue
        if current is None:
            intro.append(line)
        else:
            current["lines"].append(line)
    return {"metadata": metadata, "title": title or "Untitled artefact", "intro": intro, "sections": sections}


def _load_manifest(path: Path) -> dict[str, Any]:
    if not path.is_file():
        raise FileNotFoundError(f"No such template manifest: {path}")
    payload = yaml.safe_load(path.read_text(encoding="utf-8")) or {}
    if not isinstance(payload, dict):
        raise ValueError("Template manifest must be a YAML mapping.")
    return payload


def _section_lookup(markdown_doc: dict[str, Any]) -> dict[str, dict[str, Any]]:
    lookup: dict[str, dict[str, Any]] = {}
    for section in markdown_doc["sections"]:
        title = str(section.get("title") or "").strip()
        if title:
            lookup[title.lower()] = section
    return lookup


def _manifest_required_sections(manifest: dict[str, Any]) -> list[str]:
    value = manifest.get("required_sections") or []
    if not isinstance(value, list):
        return []
    return [str(item).strip() for item in value if str(item).strip()]


def _content_lines(section: dict[str, Any]) -> list[str]:
    return [line.rstrip() for line in section.get("lines", [])]


def _add_markdown_lines_to_docx(doc: Document, lines: list[str]) -> None:
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped[2:].strip(), style="List Bullet")
        elif re.match(r"^\d+\.\s+", stripped):
            doc.add_paragraph(re.sub(r"^\d+\.\s+", "", stripped), style="List Number")
        else:
            doc.add_paragraph(stripped)


def _resolve_docx_template(manifest_path: Path, manifest: dict[str, Any]) -> Path | None:
    template_file = str(manifest.get("template_file") or "").strip()
    if not template_file:
        return None
    template_path = (manifest_path.parent / template_file).resolve()
    root = ROOT.resolve()
    if template_path != root and root not in template_path.parents:
        raise ValueError("Template file escapes workspace root.")
    if not template_path.exists():
        return None
    if template_path.suffix.lower() != ".docx":
        raise ValueError("DOCX template_file must point to a .docx file.")
    return template_path


def _export_report(
    *,
    source_path: Path,
    manifest_path: Path,
    output_path: Path,
    output_type: str,
    warnings: list[str],
    required_sections: list[str],
    found_sections: list[str],
) -> dict[str, Any]:
    return {
        "schema_version": "document_export_report_v1",
        "source_path": _relative(source_path),
        "template_manifest_path": _relative(manifest_path),
        "output_path": _relative(output_path),
        "output_type": output_type,
        "required_sections": required_sections,
        "found_sections": found_sections,
        "warnings": warnings,
    }


@mcp.tool()
def inspect_document_template_manifest(path: str) -> dict[str, Any]:
    """Read a document export template manifest."""
    log_event(f"inspect_document_template_manifest path={path}")
    manifest_path = _safe_path(path)
    manifest = _load_manifest(manifest_path)
    return {
        "path": _relative(manifest_path),
        "template_id": manifest.get("template_id"),
        "template_name": manifest.get("template_name"),
        "output_type": manifest.get("output_type"),
        "required_sections": _manifest_required_sections(manifest),
        "template_file": manifest.get("template_file"),
        "description": manifest.get("description"),
    }


@mcp.tool()
def render_docx_from_markdown(source_path: str, template_manifest_path: str, output_path: str) -> dict[str, Any]:
    """Render a Markdown artefact into a DOCX file using a template manifest."""
    log_event(
        "render_docx_from_markdown "
        f"source_path={source_path} template_manifest_path={template_manifest_path} output_path={output_path}"
    )
    source = _safe_path(source_path)
    manifest_file = _safe_path(template_manifest_path)
    output = _safe_path(output_path)
    _enforce_output_path(output, ".docx")

    manifest = _load_manifest(manifest_file)
    if str(manifest.get("output_type") or "").lower() not in {"docx", "word"}:
        raise ValueError("Template manifest output_type must be docx for DOCX rendering.")
    markdown_doc = _parse_markdown(source.read_text(encoding="utf-8"))
    lookup = _section_lookup(markdown_doc)
    required = _manifest_required_sections(manifest)
    found = [section["title"] for section in markdown_doc["sections"]]
    warnings = [f"Missing required section: {name}" for name in required if name.lower() not in lookup]

    template_path = _resolve_docx_template(manifest_file, manifest)
    doc = Document(str(template_path)) if template_path is not None else Document()
    if template_path is None and manifest.get("template_file"):
        warnings.append("Binary DOCX template_file was declared but not found; generated from manifest structure only.")

    doc.add_heading(str(markdown_doc["title"]), level=0)
    for raw in markdown_doc["intro"]:
        if raw.strip():
            doc.add_paragraph(raw.strip())
    for section in markdown_doc["sections"]:
        level = max(1, min(int(section.get("level") or 2), 4))
        title = str(section.get("title") or "").strip()
        if title:
            doc.add_heading(title, level=level)
        _add_markdown_lines_to_docx(doc, _content_lines(section))

    output.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output))
    return _export_report(
        source_path=source,
        manifest_path=manifest_file,
        output_path=output,
        output_type="docx",
        warnings=warnings,
        required_sections=required,
        found_sections=found,
    )


@mcp.tool()
def render_pptx_from_markdown(source_path: str, template_manifest_path: str, output_path: str) -> dict[str, Any]:
    """Render a Markdown artefact into a simple PPTX deck using a template manifest."""
    log_event(
        "render_pptx_from_markdown "
        f"source_path={source_path} template_manifest_path={template_manifest_path} output_path={output_path}"
    )
    source = _safe_path(source_path)
    manifest_file = _safe_path(template_manifest_path)
    output = _safe_path(output_path)
    _enforce_output_path(output, ".pptx")

    manifest = _load_manifest(manifest_file)
    if str(manifest.get("output_type") or "").lower() not in {"pptx", "powerpoint"}:
        raise ValueError("Template manifest output_type must be pptx for PPTX rendering.")
    markdown_doc = _parse_markdown(source.read_text(encoding="utf-8"))
    lookup = _section_lookup(markdown_doc)
    required = _manifest_required_sections(manifest)
    found = [section["title"] for section in markdown_doc["sections"]]
    warnings = [f"Missing required section: {name}" for name in required if name.lower() not in lookup]

    presentation = Presentation()
    title_slide_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = str(markdown_doc["title"])
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = str(manifest.get("template_name") or "Generated from Markdown")

    bullet_layout = presentation.slide_layouts[1]
    max_bullets = int(manifest.get("max_bullets_per_slide") or 6)
    for section in markdown_doc["sections"]:
        title = str(section.get("title") or "").strip()
        if not title:
            continue
        lines = [line.strip().lstrip("-* ").strip() for line in _content_lines(section) if line.strip()]
        if not lines:
            lines = ["No content supplied."]
        slide = presentation.slides.add_slide(bullet_layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        for idx, line in enumerate(lines[:max_bullets]):
            para = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            para.text = line
            para.level = 0
        if len(lines) > max_bullets:
            para = body.add_paragraph()
            para.text = f"{len(lines) - max_bullets} additional bullet(s) omitted; see source Markdown."
            para.level = 0
            warnings.append(f"Section '{title}' was truncated for slide readability.")

    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output))
    return _export_report(
        source_path=source,
        manifest_path=manifest_file,
        output_path=output,
        output_type="pptx",
        warnings=warnings,
        required_sections=required,
        found_sections=found,
    )


if __name__ == "__main__":
    _configure_mcp_logging()
    log_event(f"document_export_server_started root={ROOT}")
    mcp.run()
